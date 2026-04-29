from __future__ import annotations

import fnmatch
from dataclasses import dataclass
from pathlib import PurePosixPath
from uuid import UUID

import requests
from lattice_jit.contracts import (
    EdgeType,
    KnowledgeEdge,
    KnowledgeNode,
    NodeType,
    SnapshotResponse,
    SnapshotStatus,
)
from lattice_jit.core import generate_id, stable_hash, utcnow
from lattice_jit.storage import SourceSnapshotRecord, StorageRepository


@dataclass(slots=True)
class SharePointSnapshotService:
    """Connector that ingests documents from a SharePoint document library via Microsoft Graph API.

    Two-phase ingestion:
      1. ``create_pending_snapshot`` creates a SOURCE node and a PENDING snapshot
         record.
      2. ``continue_ingest`` re-hydrates the snapshot, authenticates via client
         credentials, enumerates files recursively, downloads & extracts text,
         then persists SECTION nodes with BELONGS_TO edges.

    ``ingest`` is a convenience wrapper that calls both phases in sequence.
    """

    repository: StorageRepository

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def ingest(
        self,
        *,
        tenant_id: UUID,
        site_url: str,
        drive_name: str,
        folder_path: str,
        file_patterns: list[str] | None = None,
        client_id: str = "",
        client_secret: str = "",
        azure_tenant_id: str = "",
    ) -> SnapshotResponse:
        """Convenience wrapper: create pending snapshot then immediately continue ingest.

        Parameters
        ----------
        tenant_id:
            Lattice-JIT tenant identifier.
        site_url:
            Full SharePoint site URL, e.g.
            ``https://contoso.sharepoint.com/sites/MySite``.
        drive_name:
            Name of the target document library (drive).
        folder_path:
            Root-relative path within the document library, e.g. ``Reports``
            or ``Shared Documents/2025``.  Pass ``""`` or ``"/"`` to list
            from the drive root.
        file_patterns:
            Optional glob patterns to restrict which files are ingested
            (e.g. ``["*.docx", "*.xlsx"]``).  When ``None`` all supported
            extensions (``.docx``, ``.xlsx``, ``.pptx``, ``.txt``, ``.md``,
            ``.csv``) are included.
        client_id:
            Azure AD app registration client ID.
        client_secret:
            Azure AD app registration client secret.
        azure_tenant_id:
            Azure AD tenant (directory) ID used for authentication.
        """
        snapshot_id = self.create_pending_snapshot(
            tenant_id=tenant_id,
            site_url=site_url,
            drive_name=drive_name,
            folder_path=folder_path,
            file_patterns=file_patterns,
        )
        return self.continue_ingest(
            snapshot_id,
            client_id=client_id,
            client_secret=client_secret,
            azure_tenant_id=azure_tenant_id,
        )

    def create_pending_snapshot(
        self,
        *,
        tenant_id: UUID,
        site_url: str,
        drive_name: str,
        folder_path: str,
        file_patterns: list[str] | None = None,
    ) -> UUID:
        """Phase 1: create a pending snapshot record and root SOURCE node.

        Returns the newly allocated ``snapshot_id``.
        """
        snapshot_id = generate_id()
        root_node = KnowledgeNode(
            tenant_id=tenant_id,
            snapshot_id=snapshot_id,
            node_type=NodeType.SOURCE,
            title=f"SharePoint Snapshot of {site_url}/{drive_name}/{folder_path}",
            source_uri=site_url,
            body_ptr=f"{drive_name}/{folder_path}",
            content_hash=stable_hash(site_url, drive_name, folder_path),
            source_confidence=1.0,
            serving_confidence=1.0,
        )

        record = SourceSnapshotRecord(
            snapshot_id=snapshot_id,
            tenant_id=tenant_id,
            repo_path=f"{site_url}|{drive_name}",
            git_ref=folder_path if folder_path and folder_path != "/" else "",
            include_globs=file_patterns or [],
            exclude_globs=[],
            status=SnapshotStatus.PENDING,
            root_node_id=root_node.node_id,
            created_at=utcnow(),
        )
        self.repository.create_source_snapshot(record)
        self.repository.upsert_nodes([root_node])
        return snapshot_id

    def continue_ingest(
        self,
        snapshot_id: UUID,
        *,
        client_id: str = "",
        client_secret: str = "",
        azure_tenant_id: str = "",
    ) -> SnapshotResponse:
        """Phase 2: authenticate, enumerate files, extract text, persist graph.

        Credentials are accepted as keyword arguments so that ``create_pending``
        (which stores site metadata only) and the actual data-fetch step can be
        decoupled.
        """
        snapshot = self.repository.get_source_snapshot(snapshot_id)
        if snapshot is None:
            raise ValueError(f"Snapshot {snapshot_id} was not found.")
        if snapshot.status == SnapshotStatus.COMPLETED and snapshot.root_node_id is not None:
            return SnapshotResponse(
                tenant_id=snapshot.tenant_id,
                snapshot_id=snapshot.snapshot_id,
                root_node_id=snapshot.root_node_id,
                status=snapshot.status,
            )
        if snapshot.root_node_id is None:
            raise ValueError(f"Snapshot {snapshot_id} is missing a root node.")

        root_node = self.repository.get_node(snapshot.root_node_id)
        if root_node is None:
            raise ValueError(
                f"Root node {snapshot.root_node_id} was not found for snapshot {snapshot_id}."
            )

        # Decode stored metadata -------------------------------------------------
        try:
            site_url, drive_name = snapshot.repo_path.split("|", 1)
        except ValueError:
            raise ValueError(f"Invalid repo_path encoding: {snapshot.repo_path!r}") from None
        folder_path = snapshot.git_ref or ""
        file_patterns = snapshot.include_globs

        # Authenticate & enumerate -----------------------------------------------
        access_token = self._get_access_token(azure_tenant_id, client_id, client_secret)
        site_id = self._get_site_id(access_token, site_url)
        drive_id = self._get_drive_id(access_token, site_id, drive_name)

        file_infos: list[dict[str, str]] = []
        self._enumerate_files(access_token, site_id, drive_id, folder_path, file_infos)

        # Download, extract & build graph ----------------------------------------
        nodes: list[KnowledgeNode] = [root_node]
        edges: list[KnowledgeEdge] = []

        for file_info in file_infos:
            file_name = file_info["name"]
            file_path = file_info["path"]

            if not self._matches_pattern(file_name, file_patterns):
                continue

            content = self._download_and_extract_text(
                file_info["download_url"], file_name
            )
            if content is None:
                continue

            section_node = KnowledgeNode(
                tenant_id=snapshot.tenant_id,
                snapshot_id=snapshot.snapshot_id,
                node_type=NodeType.SECTION,
                title=file_name,
                source_uri=f"{site_url}/{drive_name}/{file_path}",
                body_ptr=file_path,
                body_text=content,
                content_hash=stable_hash(site_url, drive_name, file_path, content),
                source_confidence=1.0,
                serving_confidence=1.0,
            )
            nodes.append(section_node)
            edges.append(
                KnowledgeEdge(
                    tenant_id=snapshot.tenant_id,
                    from_node_id=section_node.node_id,
                    to_node_id=root_node.node_id,
                    edge_type=EdgeType.BELONGS_TO,
                    evidence_spans=[{"path": file_path}],
                )
            )

        self.repository.upsert_nodes(nodes)
        self.repository.upsert_edges(edges)
        self.repository.mark_snapshot_completed(snapshot.snapshot_id, root_node.node_id)
        return SnapshotResponse(
            tenant_id=snapshot.tenant_id,
            snapshot_id=snapshot.snapshot_id,
            root_node_id=root_node.node_id,
            status=SnapshotStatus.COMPLETED,
        )

    # ------------------------------------------------------------------
    # SharePoint / Microsoft Graph API helpers
    # ------------------------------------------------------------------

    def _get_access_token(self, tenant: str, client_id: str, client_secret: str) -> str:
        """Authenticate with Azure AD via the OAuth 2.0 client-credentials flow.

        Returns an access token scoped to ``https://graph.microsoft.com/.default``.
        """
        if not tenant or not client_id or not client_secret:
            raise RuntimeError(
                "SharePoint authentication requires azure_tenant_id, client_id, "
                "and client_secret."
            )
        url = f"https://login.microsoftonline.com/{tenant}/oauth2/v2.0/token"
        resp = requests.post(
            url,
            data={
                "client_id": client_id,
                "client_secret": client_secret,
                "scope": "https://graph.microsoft.com/.default",
                "grant_type": "client_credentials",
            },
            timeout=30,
        )
        try:
            resp.raise_for_status()
        except requests.HTTPError as exc:
            raise RuntimeError(
                f"Azure AD auth failed (HTTP {resp.status_code}): {resp.text}"
            ) from exc
        return str(resp.json()["access_token"])

    def _get_site_id(self, access_token: str, site_url: str) -> str:
        """Resolve a SharePoint site URL to a Graph API ``sites/{id}``."""
        from urllib.parse import urlparse

        parsed = urlparse(site_url)
        hostname = parsed.netloc
        server_relative_path = parsed.path  # e.g. ``/sites/MySite``
        api_url = (
            f"https://graph.microsoft.com/v1.0/sites/{hostname}:{server_relative_path}"
        )
        resp = requests.get(
            api_url,
            headers={"Authorization": f"Bearer {access_token}"},
            timeout=30,
        )
        try:
            resp.raise_for_status()
        except requests.HTTPError as exc:
            raise RuntimeError(
                f"Failed to resolve site {site_url} (HTTP {resp.status_code}): "
                f"{resp.text}"
            ) from exc
        return str(resp.json()["id"])

    def _get_drive_id(self, access_token: str, site_id: str, drive_name: str) -> str:
        """Find a document library (drive) by name on the given SharePoint site."""
        api_url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives"
        resp = requests.get(
            api_url,
            headers={"Authorization": f"Bearer {access_token}"},
            timeout=30,
        )
        try:
            resp.raise_for_status()
        except requests.HTTPError as exc:
            raise RuntimeError(
                f"Failed to list drives for site {site_id} (HTTP {resp.status_code}): "
                f"{resp.text}"
            ) from exc
        drives = resp.json().get("value", [])
        for drive in drives:
            if drive.get("name") == drive_name:
                return str(drive["id"])
        raise ValueError(f"Drive {drive_name!r} not found on site {site_id}.")

    def _enumerate_files(
        self,
        access_token: str,
        site_id: str,
        drive_id: str,
        folder_path: str,
        results: list[dict[str, str]],
    ) -> None:
        """Recursively enumerate all files within *folder_path* on the drive.

        Each entry in *results* has keys ``name``, ``path``, and ``download_url``.
        Pagination via ``@odata.nextLink`` is handled transparently.
        """
        # Build the initial children endpoint
        stripped = folder_path.strip("/")
        if stripped:
            base_url = (
                f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}"
                f"/root:/{stripped}:/children"
            )
        else:
            base_url = (
                f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}"
                f"/root/children"
            )

        next_url: str | None = base_url
        while next_url:
            resp = requests.get(
                next_url,
                headers={"Authorization": f"Bearer {access_token}"},
                timeout=30,
            )
            try:
                resp.raise_for_status()
            except requests.HTTPError as exc:
                raise RuntimeError(
                    f"Failed to list children at {folder_path!r} "
                    f"(HTTP {resp.status_code}): {resp.text}"
                ) from exc

            data = resp.json()
            for item in data.get("value", []):
                if "folder" in item:
                    subfolder = f"{stripped}/{item['name']}".strip("/")
                    self._enumerate_files(
                        access_token, site_id, drive_id, subfolder, results
                    )
                elif "file" in item:
                    download_url = item.get("@microsoft.graph.downloadUrl")
                    if download_url:
                        file_path = f"{stripped}/{item['name']}".strip("/")
                        results.append(
                            {
                                "name": item["name"],
                                "path": file_path,
                                "download_url": download_url,
                            }
                        )

            next_url = data.get("@odata.nextLink")

    # ------------------------------------------------------------------
    # File matching
    # ------------------------------------------------------------------

    SUPPORTED_EXTENSIONS: frozenset[str] = frozenset(
        {".docx", ".xlsx", ".pptx", ".txt", ".md", ".csv"}
    )

    def _matches_pattern(self, file_name: str, patterns: list[str]) -> bool:
        """Return ``True`` if *file_name* matches any of the glob *patterns*.

        When *patterns* is empty every supported extension is accepted.
        """
        if not patterns:
            return PurePosixPath(file_name).suffix.lower() in self.SUPPORTED_EXTENSIONS
        return any(fnmatch.fnmatch(file_name, pat) for pat in patterns)

    # ------------------------------------------------------------------
    # Download & text extraction
    # ------------------------------------------------------------------

    def _download_and_extract_text(self, download_url: str, file_name: str) -> str | None:
        """Download file content and extract human-readable text.

        Returns ``None`` when the file type is not supported or extraction
        fails silently.
        """
        resp = requests.get(download_url, timeout=120)
        try:
            resp.raise_for_status()
        except requests.HTTPError:
            return None
        return self._extract_text(resp.content, file_name)

    def _extract_text(self, content: bytes, file_name: str) -> str | None:
        """Dispatch to the appropriate extractor based on file extension."""
        ext = PurePosixPath(file_name).suffix.lower()

        if ext in (".txt", ".md", ".csv"):
            return self._decode_text(content)
        if ext == ".docx":
            return self._extract_docx(content)
        if ext == ".xlsx":
            return self._extract_xlsx(content)
        if ext == ".pptx":
            return self._extract_pptx(content)
        return None

    # -- plain-text -----------------------------------------------------

    @staticmethod
    def _decode_text(content: bytes) -> str:
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            return content.decode("utf-8", errors="replace")

    # -- docx -----------------------------------------------------------

    @staticmethod
    def _extract_docx(content: bytes) -> str:
        try:
            from docx import Document
        except ImportError as exc:
            raise RuntimeError(
                "python-docx is not installed. Install it with: pip install python-docx"
            ) from exc

        import io as _io
        doc = Document(_io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text]
        return "\n".join(paragraphs)

    # -- xlsx -----------------------------------------------------------

    @staticmethod
    def _extract_xlsx(content: bytes) -> str:
        try:
            import openpyxl
        except ImportError as exc:
            raise RuntimeError(
                "openpyxl is not installed. Install it with: pip install openpyxl"
            ) from exc

        wb = openpyxl.load_workbook(content, read_only=True, data_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                rows.append("\t".join(cells))
            parts.append(f"--- {sheet_name} ---\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)

    # -- pptx -----------------------------------------------------------

    @staticmethod
    def _extract_pptx(content: bytes) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError(
                "python-pptx is not installed. Install it with: pip install python-pptx"
            ) from exc

        import io as _io2
        prs = Presentation(_io2.BytesIO(content))
        parts: list[str] = []
        for slide in prs.slides:
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text)
            if slide_texts:
                parts.append("\n".join(slide_texts))
        return "\n\n".join(parts)
